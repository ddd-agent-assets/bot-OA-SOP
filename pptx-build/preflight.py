# 程序化预检：文本溢出启发式 + 越界检查（standalone）
# 用法：python preflight.py [pptx路径或目录 ...]  默认 ./slides/*.pptx
# 判读口径：与人工标杆同款的误报可放过——单行框 need~0.70 vs has~0.50 是行高估算保守，非真溢出
import os, re, glob, unicodedata, sys
from pptx import Presentation
from pptx.util import Emu

args = sys.argv[1:] or [os.path.join(os.getcwd(), "slides")]
files = []
for a in args:
    files += sorted(glob.glob(os.path.join(a, "*.pptx"))) if os.path.isdir(a) else [a]

def char_w(ch, pt):
    if unicodedata.east_asian_width(ch) in "WF":
        return pt / 72.0
    if ch in "iljI.,:;'|! ":
        return pt / 72.0 * 0.3
    return pt / 72.0 * 0.55

issues = []
for f in files:
    deck = os.path.basename(f)
    prs = Presentation(f)
    for si, slide in enumerate(prs.slides, 1):
        def walk(shapes):
            for sp in shapes:
                if sp.shape_type == 6:
                    walk(sp.shapes); continue
                x, y = Emu(sp.left).inches, Emu(sp.top).inches
                w, h = Emu(sp.width).inches, Emu(sp.height).inches
                if x < -0.7 or y < -0.7 or x + w > 20.05 or y + h > 11.3:
                    issues.append(f"{deck} s{si}: OUT-OF-CANVAS {sp.shape_id} ({x:.2f},{y:.2f},{w:.2f}x{h:.2f})")
                if not getattr(sp, "has_text_frame", False) or not sp.has_text_frame:
                    continue
                tf = sp.text_frame
                total_lines = 0
                max_pt = 0
                for para in tf.paragraphs:
                    txt = "".join(r.text for r in para.runs)
                    if not txt.strip():
                        total_lines += 0.6
                        continue
                    pt = max((r.font.size.pt if r.font.size else 0) for r in para.runs) or 18
                    max_pt = max(max_pt, pt)
                    line_w = sum(char_w(c, pt) for c in txt)
                    total_lines += max(1, -(-line_w // max(w - 0.1, 0.5)))
                need_h = total_lines * max_pt / 72.0 * 1.35 + 0.1
                if need_h > h * 1.18 and total_lines >= 1:
                    t = tf.text[:30].replace("\n", "|")
                    issues.append(f"{deck} s{si}: OVERFLOW? id={sp.shape_id} need~{need_h:.2f}in has={h:.2f}in text={t!r}")
        walk(slide.shapes)

print(f"{len(issues)} potential issues")
for i in issues:
    print(" -", i)
